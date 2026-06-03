"""
═══════════════════════════════════════════════════════════════════════════════
SlideGenerator — AI-Powered Brand Deck Builder
═══════════════════════════════════════════════════════════════════════════════
Sinh structured JSON slides cho:
  - Brand Guideline (8-10 slides)
  - Marketing Proposal (6-8 slides) 
  - Pitch Deck (8-10 slides)
  
Output: Array of slide objects ready for frontend rendering & PPTX export
═══════════════════════════════════════════════════════════════════════════════
"""

import os
import json
import uuid
from typing import Dict, Any, List, Optional
from io import BytesIO


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE TEMPLATES — Enterprise-Grade Prompts
# ═══════════════════════════════════════════════════════════════════════════════

BRAND_GUIDELINE_PROMPT = """Bạn là Brand Identity Director tại Pentagram — agency thiết kế hàng đầu thế giới.

Từ Brand DNA dưới đây, hãy sinh ra một bộ Brand Guideline hoàn chỉnh gồm 8 slides.

<BRAND_DNA>
{brand_dna}
</BRAND_DNA>

<BUSINESS_CONTEXT>
{business_context}
</BUSINESS_CONTEXT>

MỖI SLIDE phải có cấu trúc JSON chính xác:
{{
  "slide_id": "unique-id",
  "slide_number": 1,
  "layout": "cover|section|content|split|gallery|palette|typography|rules",
  "background": {{
    "type": "solid|gradient",
    "color": "#hex hoặc linear-gradient(...)",
    "dark_mode": true/false
  }},
  "elements": [
    {{
      "type": "heading|subheading|body|bullet_list|color_swatch|font_sample|icon_text|image_placeholder|divider|badge",
      "content": "Nội dung text",
      "style": {{
        "color": "#hex",
        "fontSize": 48,
        "fontWeight": "bold|normal|light",
        "textAlign": "left|center|right",
        "x": 0, "y": 0, "w": 100, "h": 20
      }}
    }}
  ]
}}

Các vị trí x,y,w,h tính theo PHẦN TRĂM (0-100) so với kích thước slide.

═══ 8 SLIDES BẮT BUỘC ═══

1. **Cover** (layout: "cover"): Tên thương hiệu lớn + tagline + năm. Background gradient sang trọng dùng màu brand.
2. **Brand Story** (layout: "section"): Mission & Vision statement. Ngắn gọn, truyền cảm hứng.
3. **Brand Values** (layout: "content"): 3-4 giá trị cốt lõi, mỗi giá trị có icon text + mô tả 1 câu.
4. **Logo Usage** (layout: "split"): Quy tắc sử dụng logo — clear space, minimum size, do/don't.
5. **Color Palette** (layout: "palette"): Primary + Secondary + Accent colors (tối thiểu 5 mã HEX). Mỗi color swatch kèm tên và HEX code.
6. **Typography** (layout: "typography"): Heading font + Body font + hierarchy (H1 > H2 > Body > Caption). Kèm sample text.
7. **Tone of Voice** (layout: "content"): Giọng điệu thương hiệu — We are / We are not. Kèm 2-3 ví dụ copy mẫu.
8. **Brand Application** (layout: "gallery"): Mô tả 3-4 mockup ứng dụng (namecard, uniform, social media post, packaging).

═══ TIÊU CHUẨN THIẾT KẾ ═══
- Màu sắc PHẢI lấy từ Brand DNA (nếu có), hoặc sinh palette phù hợp ngành.
- Typography: Đề xuất Google Fonts thực tế (Inter, Montserrat, Playfair Display, Poppins...)
- Mọi text phải CHÍNH XÁC cho doanh nghiệp cụ thể, KHÔNG dùng lorem ipsum.
- Thiết kế chuẩn enterprise, đủ để CEO/CMO dùng trình bày trong meeting.

Trả về ĐÚNG JSON array: [slide1, slide2, ...]. KHÔNG thêm markdown backticks.
"""

PITCH_DECK_PROMPT = """Bạn là Startup Advisor kiêm Pitch Coach đã mentor 200+ startups gọi vốn thành công tại Đông Nam Á.

Từ Brand DNA và Business Context dưới đây, hãy sinh ra một Pitch Deck gồm 8 slides.

<BRAND_DNA>
{brand_dna}
</BRAND_DNA>

<BUSINESS_CONTEXT>
{business_context}
</BUSINESS_CONTEXT>

MỖI SLIDE phải có cấu trúc JSON chính xác giống format đã mô tả.

═══ 8 SLIDES BẮT BUỘC ═══

1. **Cover**: Tên công ty + tagline + logo placeholder.
2. **Problem**: Pain point thị trường — dùng số liệu cụ thể VN.
3. **Solution**: Giải pháp/sản phẩm — USP rõ ràng.
4. **Market Size**: TAM/SAM/SOM bằng số VND/tỷ. Circle diagram concept.
5. **Business Model**: Cách kiếm tiền — Revenue streams.
6. **Traction**: Milestone đã đạt hoặc roadmap 12 tháng.
7. **Team**: Founder/Co-founder profiles (placeholder).
8. **The Ask**: Số tiền cần gọi + cách sử dụng (pie chart concept).

═══ TIÊU CHUẨN ═══
- Thiết kế tối giản, hiện đại (dark theme hoặc light clean).
- Mỗi slide chỉ 1 key message — KHÔNG nhồi nhét thông tin.
- Số liệu phải thực tế cho thị trường Việt Nam.

Trả về ĐÚNG JSON array: [slide1, slide2, ...]. KHÔNG thêm markdown backticks.
"""

PROPOSAL_PROMPT = """Bạn là Strategy Director tại Ogilvy Vietnam với 15 năm kinh nghiệm viết proposal cho các brand lớn.

Từ Brand DNA và Business Context, sinh ra Marketing Proposal gồm 7 slides.

<BRAND_DNA>
{brand_dna}
</BRAND_DNA>

<BUSINESS_CONTEXT>
{business_context}
</BUSINESS_CONTEXT>

═══ 7 SLIDES ═══

1. **Cover**: Tiêu đề proposal + tên khách hàng + ngày.
2. **Executive Summary**: Tóm tắt 3-4 bullet points key takeaways.
3. **Situation Analysis**: Thị trường + insight khách hàng + cơ hội.
4. **Strategy**: Big idea + campaign concept + key message.
5. **Execution Plan**: Timeline 3 tháng + kênh triển khai.
6. **Budget**: Bảng phân bổ ngân sách theo kênh (VND).
7. **Why Us / CTA**: Lý do chọn + call to action.

Trả về ĐÚNG JSON array. KHÔNG thêm markdown backticks.
"""


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE GENERATOR CLASS
# ═══════════════════════════════════════════════════════════════════════════════

class SlideGenerator:
    """AI-powered slide deck generator for Brand Guidelines, Proposals, Pitch Decks."""
    
    TEMPLATE_PROMPTS = {
        "brand_guideline": BRAND_GUIDELINE_PROMPT,
        "pitch_deck": PITCH_DECK_PROMPT,
        "proposal": PROPOSAL_PROMPT,
    }
    
    def __init__(self):
        api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY")
        self.llm = None
        if api_key:
            from langchain_google_genai import ChatGoogleGenerativeAI
            self.llm = ChatGoogleGenerativeAI(
                model="gemini-2.5-flash",
                temperature=0.4,
                api_key=api_key,
                max_retries=2,
                timeout=120.0
            )
    
    def generate_slides(
        self, 
        template_type: str, 
        brand_dna: Dict[str, Any] = None, 
        business_context: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Generate a complete slide deck from Brand DNA.
        
        Args:
            template_type: "brand_guideline" | "pitch_deck" | "proposal"
            brand_dna: Brand DNA dict from onboarding
            business_context: Additional business context
            
        Returns:
            {"status": "success", "data": {"slides": [...], "metadata": {...}}}
        """
        if not self.llm:
            return self._generate_mock_slides(template_type, brand_dna, business_context)
        
        prompt_template = self.TEMPLATE_PROMPTS.get(template_type)
        if not prompt_template:
            return {"status": "error", "message": f"Unknown template: {template_type}"}
        
        dna_str = json.dumps(brand_dna or {}, ensure_ascii=False, indent=2)
        ctx_str = json.dumps(business_context or {}, ensure_ascii=False, indent=2)
        
        prompt = prompt_template.format(brand_dna=dna_str, business_context=ctx_str)
        
        try:
            from langchain_core.messages import HumanMessage
            response = self.llm.invoke([HumanMessage(content=prompt)])
            raw = response.content.strip()
            
            # Clean markdown artifacts
            if raw.startswith("```"):
                raw = raw.split("```", 2)[1]
                if raw.startswith("json"):
                    raw = raw[4:]
                raw = raw.strip()
            if raw.endswith("```"):
                raw = raw[:-3].strip()
            
            slides = json.loads(raw)
            
            # Ensure IDs
            for i, slide in enumerate(slides):
                if not slide.get("slide_id"):
                    slide["slide_id"] = f"slide-{uuid.uuid4().hex[:8]}"
                slide["slide_number"] = i + 1
            
            return {
                "status": "success",
                "data": {
                    "slides": slides,
                    "metadata": {
                        "template_type": template_type,
                        "slide_count": len(slides),
                        "brand_name": (brand_dna or {}).get("brand_name", ""),
                    }
                }
            }
        except json.JSONDecodeError as e:
            print(f"⚠️ [SLIDE GEN] JSON parse error: {e}")
            return self._generate_mock_slides(template_type, brand_dna, business_context)
        except Exception as e:
            print(f"❌ [SLIDE GEN] Error: {e}")
            return {"status": "error", "message": str(e)}
    
    def _generate_mock_slides(
        self, 
        template_type: str, 
        brand_dna: Dict[str, Any] = None,
        business_context: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """Generate mock slides for demo/dev without API key."""
        brand_name = (brand_dna or {}).get("brand_name") or (business_context or {}).get("company_name") or "BrandFlow"
        industry = (business_context or {}).get("industry") or "Marketing"
        tagline = (brand_dna or {}).get("positioning") or "Elevate Your Brand"
        tone = (brand_dna or {}).get("tone_of_voice") or "Chuyên nghiệp & Hiện đại"
        
        primary_color = "#0F172A"
        accent_color = "#06B6D4"
        
        slides = [
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 1,
                "layout": "cover",
                "background": {"type": "gradient", "color": f"linear-gradient(135deg, {primary_color} 0%, #1E293B 100%)", "dark_mode": True},
                "elements": [
                    {"type": "badge", "content": "BRAND GUIDELINES 2025", "style": {"color": accent_color, "fontSize": 12, "fontWeight": "bold", "textAlign": "center", "x": 30, "y": 25, "w": 40, "h": 5}},
                    {"type": "heading", "content": brand_name.upper(), "style": {"color": "#FFFFFF", "fontSize": 64, "fontWeight": "bold", "textAlign": "center", "x": 10, "y": 35, "w": 80, "h": 15}},
                    {"type": "subheading", "content": tagline, "style": {"color": "#94A3B8", "fontSize": 20, "fontWeight": "light", "textAlign": "center", "x": 20, "y": 52, "w": 60, "h": 8}},
                    {"type": "divider", "content": "", "style": {"color": accent_color, "fontSize": 1, "fontWeight": "normal", "textAlign": "center", "x": 40, "y": 64, "w": 20, "h": 1}},
                    {"type": "body", "content": f"Ngành: {industry}", "style": {"color": "#64748B", "fontSize": 14, "fontWeight": "normal", "textAlign": "center", "x": 30, "y": 70, "w": 40, "h": 5}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 2,
                "layout": "section",
                "background": {"type": "solid", "color": "#FFFFFF", "dark_mode": False},
                "elements": [
                    {"type": "badge", "content": "01 — BRAND STORY", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 15, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Mission & Vision", "style": {"color": primary_color, "fontSize": 48, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 25, "w": 60, "h": 12}},
                    {"type": "body", "content": f"Chúng tôi tin rằng mỗi doanh nghiệp {industry} đều xứng đáng có một thương hiệu mạnh mẽ. {brand_name} ra đời với sứ mệnh giúp các SME Việt Nam xây dựng bản sắc riêng, cạnh tranh bình đẳng với các tập đoàn lớn thông qua chiến lược thương hiệu thông minh và sáng tạo.", "style": {"color": "#475569", "fontSize": 16, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 42, "w": 55, "h": 25}},
                    {"type": "body", "content": f"Tầm nhìn: Trở thành đối tác Brand Strategy #1 cho SME {industry} tại Việt Nam vào năm 2027.", "style": {"color": primary_color, "fontSize": 18, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 72, "w": 55, "h": 10}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 3,
                "layout": "content",
                "background": {"type": "solid", "color": "#F8FAFC", "dark_mode": False},
                "elements": [
                    {"type": "badge", "content": "02 — BRAND VALUES", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 10, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Giá trị Cốt lõi", "style": {"color": primary_color, "fontSize": 40, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 18, "w": 50, "h": 10}},
                    {"type": "icon_text", "content": "🎯 Chính xác — Mọi quyết định dựa trên dữ liệu, không phỏng đoán.", "style": {"color": "#334155", "fontSize": 16, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 35, "w": 80, "h": 8}},
                    {"type": "icon_text", "content": "💡 Sáng tạo — Luôn tìm cách mới để giải bài toán cũ.", "style": {"color": "#334155", "fontSize": 16, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 47, "w": 80, "h": 8}},
                    {"type": "icon_text", "content": "🤝 Đồng hành — Không chỉ giao output, mà đi cùng đến kết quả.", "style": {"color": "#334155", "fontSize": 16, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 59, "w": 80, "h": 8}},
                    {"type": "icon_text", "content": "⚡ Tốc độ — Time-to-market nhanh gấp 3x so với agency truyền thống.", "style": {"color": "#334155", "fontSize": 16, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 71, "w": 80, "h": 8}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 4,
                "layout": "split",
                "background": {"type": "solid", "color": "#FFFFFF", "dark_mode": False},
                "elements": [
                    {"type": "badge", "content": "03 — LOGO USAGE", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 10, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Quy tắc Sử dụng Logo", "style": {"color": primary_color, "fontSize": 36, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 18, "w": 45, "h": 10}},
                    {"type": "image_placeholder", "content": f"[Logo {brand_name}]", "style": {"color": "#E2E8F0", "fontSize": 14, "fontWeight": "normal", "textAlign": "center", "x": 55, "y": 15, "w": 35, "h": 35}},
                    {"type": "bullet_list", "content": "✅ Clear space: Tối thiểu 1x chiều cao logo\n✅ Kích thước tối thiểu: 24px (digital), 10mm (print)\n✅ Luôn dùng file vector (SVG/AI) cho in ấn\n❌ KHÔNG kéo méo, xoay, hoặc đổ bóng logo\n❌ KHÔNG đặt logo trên nền rối hoặc ảnh nhiều chi tiết\n❌ KHÔNG thay đổi màu sắc ngoài quy định", "style": {"color": "#475569", "fontSize": 14, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 35, "w": 42, "h": 50}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 5,
                "layout": "palette",
                "background": {"type": "solid", "color": "#FFFFFF", "dark_mode": False},
                "elements": [
                    {"type": "badge", "content": "04 — COLOR PALETTE", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 10, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Bảng Màu Thương hiệu", "style": {"color": primary_color, "fontSize": 36, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 18, "w": 50, "h": 10}},
                    {"type": "color_swatch", "content": f"Primary\n{primary_color}", "style": {"color": primary_color, "fontSize": 12, "fontWeight": "bold", "textAlign": "center", "x": 8, "y": 35, "w": 16, "h": 35}},
                    {"type": "color_swatch", "content": f"Accent\n{accent_color}", "style": {"color": accent_color, "fontSize": 12, "fontWeight": "bold", "textAlign": "center", "x": 26, "y": 35, "w": 16, "h": 35}},
                    {"type": "color_swatch", "content": "Secondary\n#1E293B", "style": {"color": "#1E293B", "fontSize": 12, "fontWeight": "bold", "textAlign": "center", "x": 44, "y": 35, "w": 16, "h": 35}},
                    {"type": "color_swatch", "content": "Light\n#F1F5F9", "style": {"color": "#F1F5F9", "fontSize": 12, "fontWeight": "bold", "textAlign": "center", "x": 62, "y": 35, "w": 16, "h": 35}},
                    {"type": "color_swatch", "content": "Warm\n#F59E0B", "style": {"color": "#F59E0B", "fontSize": 12, "fontWeight": "bold", "textAlign": "center", "x": 80, "y": 35, "w": 16, "h": 35}},
                    {"type": "body", "content": "Luôn đảm bảo độ tương phản WCAG AA (tối thiểu 4.5:1) khi kết hợp text trên nền màu.", "style": {"color": "#64748B", "fontSize": 13, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 78, "w": 80, "h": 8}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 6,
                "layout": "typography",
                "background": {"type": "solid", "color": "#F8FAFC", "dark_mode": False},
                "elements": [
                    {"type": "badge", "content": "05 — TYPOGRAPHY", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 10, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Hệ thống Chữ", "style": {"color": primary_color, "fontSize": 36, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 18, "w": 45, "h": 10}},
                    {"type": "font_sample", "content": "Inter — Heading Font\nAaBbCcDd 0123456789", "style": {"color": primary_color, "fontSize": 32, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 35, "w": 80, "h": 15}},
                    {"type": "font_sample", "content": "Inter — Body Font\nViệt Nam là thị trường năng động bậc nhất Đông Nam Á.", "style": {"color": "#475569", "fontSize": 18, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 55, "w": 80, "h": 12}},
                    {"type": "body", "content": "H1: 48px Bold | H2: 36px Semibold | H3: 24px Medium | Body: 16px Regular | Caption: 12px Light", "style": {"color": "#94A3B8", "fontSize": 13, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 75, "w": 80, "h": 8}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 7,
                "layout": "content",
                "background": {"type": "solid", "color": "#FFFFFF", "dark_mode": False},
                "elements": [
                    {"type": "badge", "content": "06 — TONE OF VOICE", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 10, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Giọng điệu Thương hiệu", "style": {"color": primary_color, "fontSize": 36, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 18, "w": 60, "h": 10}},
                    {"type": "body", "content": f"Giọng điệu: {tone}", "style": {"color": accent_color, "fontSize": 18, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 32, "w": 80, "h": 6}},
                    {"type": "icon_text", "content": "✅ Chúng tôi là: Thẳng thắn, Có dữ liệu, Đi thẳng vào trọng tâm, Thân thiện nhưng chuyên nghiệp.", "style": {"color": "#334155", "fontSize": 15, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 42, "w": 80, "h": 8}},
                    {"type": "icon_text", "content": "❌ Chúng tôi KHÔNG: Sáo rỗng, Hứa suông không số liệu, Dùng jargon khó hiểu, Quá formal/robot.", "style": {"color": "#334155", "fontSize": 15, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 54, "w": 80, "h": 8}},
                    {"type": "body", "content": f"Ví dụ copy: \"{brand_name} giúp bạn tăng 3x doanh thu trong 90 ngày — không phải bằng quảng cáo spam, mà bằng chiến lược thương hiệu đúng người, đúng thời điểm.\"", "style": {"color": "#64748B", "fontSize": 14, "fontWeight": "normal", "textAlign": "left", "x": 10, "y": 68, "w": 80, "h": 15}},
                ]
            },
            {
                "slide_id": f"slide-{uuid.uuid4().hex[:8]}",
                "slide_number": 8,
                "layout": "gallery",
                "background": {"type": "gradient", "color": f"linear-gradient(135deg, {primary_color} 0%, #1E293B 100%)", "dark_mode": True},
                "elements": [
                    {"type": "badge", "content": "07 — BRAND APPLICATION", "style": {"color": accent_color, "fontSize": 11, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 10, "w": 30, "h": 4}},
                    {"type": "heading", "content": "Ứng dụng Thương hiệu", "style": {"color": "#FFFFFF", "fontSize": 36, "fontWeight": "bold", "textAlign": "left", "x": 10, "y": 18, "w": 60, "h": 10}},
                    {"type": "image_placeholder", "content": "🪪 Name Card", "style": {"color": "#334155", "fontSize": 14, "fontWeight": "bold", "textAlign": "center", "x": 8, "y": 35, "w": 20, "h": 25}},
                    {"type": "image_placeholder", "content": "📱 Social Media", "style": {"color": "#334155", "fontSize": 14, "fontWeight": "bold", "textAlign": "center", "x": 30, "y": 35, "w": 20, "h": 25}},
                    {"type": "image_placeholder", "content": "👕 Uniform", "style": {"color": "#334155", "fontSize": 14, "fontWeight": "bold", "textAlign": "center", "x": 52, "y": 35, "w": 20, "h": 25}},
                    {"type": "image_placeholder", "content": "📦 Packaging", "style": {"color": "#334155", "fontSize": 14, "fontWeight": "bold", "textAlign": "center", "x": 74, "y": 35, "w": 20, "h": 25}},
                    {"type": "body", "content": f"© 2025 {brand_name}. All rights reserved. Brand Guidelines v1.0", "style": {"color": "#64748B", "fontSize": 12, "fontWeight": "normal", "textAlign": "center", "x": 20, "y": 85, "w": 60, "h": 5}},
                ]
            },
        ]
        
        return {
            "status": "success",
            "data": {
                "slides": slides,
                "metadata": {
                    "template_type": template_type,
                    "slide_count": len(slides),
                    "brand_name": brand_name,
                    "mock_mode": not bool(self.llm),
                }
            }
        }


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX EXPORTER
# ═══════════════════════════════════════════════════════════════════════════════

def export_slides_to_pptx(slides: List[Dict], brand_name: str = "Brand") -> BytesIO:
    """
    Export slide JSON array to a PowerPoint (.pptx) file.
    Returns BytesIO buffer ready for streaming.
    """
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches, Pt, Emu  # type: ignore[import-untyped]
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    
    def hex_to_rgb(hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip("#")
        if len(hex_color) != 6:
            return RGBColor(15, 23, 42)
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    
    for slide_data in slides:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        bg = slide_data.get("background", {})
        bg_color = bg.get("color", "#FFFFFF")
        if not bg_color.startswith("linear"):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(bg_color)
        else:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb("#0F172A")
        
        # Elements
        for elem in slide_data.get("elements", []):
            style = elem.get("style", {})
            x_pct = style.get("x", 0) / 100
            y_pct = style.get("y", 0) / 100
            w_pct = style.get("w", 50) / 100
            h_pct = style.get("h", 10) / 100
            
            left = int(slide_w * x_pct)
            top = int(slide_h * y_pct)
            width = int(slide_w * w_pct)
            height = int(slide_h * h_pct)
            
            content = elem.get("content", "")
            elem_type = elem.get("type", "body")
            
            if elem_type == "divider":
                from pptx.util import Pt as PtUtil  # type: ignore[import-untyped]
                shape = slide.shapes.add_shape(
                    1, left, top, width, Emu(Pt(2))
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(style.get("color", "#06B6D4"))
                shape.line.fill.background()
                continue
            
            if elem_type == "color_swatch":
                shape = slide.shapes.add_shape(1, left, top, width, height)
                swatch_color = style.get("color", "#CCCCCC")
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(swatch_color)
                shape.line.fill.background()
                
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                continue
            
            if elem_type == "image_placeholder":
                shape = slide.shapes.add_shape(1, left, top, width, height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb("#E2E8F0")
                shape.line.fill.background()
                
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(12)
                p.font.color.rgb = hex_to_rgb("#475569")
                p.alignment = PP_ALIGN.CENTER
                continue
            
            # Text elements
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            font_size = style.get("fontSize", 16)
            font_color = style.get("color", "#000000")
            font_weight = style.get("fontWeight", "normal")
            text_align = align_map.get(style.get("textAlign", "left"), PP_ALIGN.LEFT)
            
            # Handle multi-line content
            lines = content.split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(min(font_size, 48))
                p.font.color.rgb = hex_to_rgb(font_color)
                p.font.bold = font_weight == "bold"
                p.alignment = text_align
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
